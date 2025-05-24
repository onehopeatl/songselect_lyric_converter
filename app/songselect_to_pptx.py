
import re
import json
import os
import logging
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# INFO logging level is default
# Use DEBUG to get additional info for troubleshooting
logging.basicConfig(level=logging.INFO)

# SET GLOBAL VARS
#PRS = Presentation()

def parse_lyrics_file(file_path, ccli_lic_num):
    """
    Parse a CCLI Song Select txt lyrics file to sections and metadata
    
    Args:
        file_path (str): Path to the lyrics file
        ccl_lic_num (str): CCL License Number
        
    Returns:
        dict: A dictionary containing the parsed song structure
    """
    # Read the file content
    with open(file_path, 'r', encoding='utf-8') as file:
        content = file.read()
    
    # Initialize the song structure
    song = {
        'title': '',
        'sections': [],
        'metadata': {}
    }
    
    # Extract the title (first non-empty line)
    lines = content.split('\n')
    for line in lines:
        if line.strip():
            song['title'] = line.strip()
            break
    
    # Define a list of potential section headings for regex matching
    # \d* matches numbers after the heading name
    heading_list = [
        'Verse',
        'Verse \d*',
        'Chorus',
        'Chorus \d*',
        'Pre-Chorus',
        'Pre-Chorus \d*',
        'Bridge',
        'Bridge \d*',
        'Tag',
        'Tag \d*',
        'Intro',
        'Intro \d*',
        'Outro',
        'Outro \d*',
        'Interlude',
        'Interlude \d*',
        'Vamp',
        'Vamp \d*',
        'Ending',
        'Ending \d*'
    ]
    # Regex to identify section headers
    section_pattern = re.compile(r'^(%s)' % '|'.join(heading_list), re.IGNORECASE)
    logging.debug(f'Regex pattern for matching sections... {section_pattern}')
    
    # Regex to identify the start of the metadata section
    metadata_pattern = re.compile(r'^(?:\s*$)|^(?:.*(?:CCLI|©|Copyright|\(C\)))|^(?:[A-Za-z]+(?: [A-Za-z]+)+(?:,|;))', re.IGNORECASE)
    logging.debug(f'Regex pattern for matching metadata section... {metadata_pattern}')
    
    current_section = None
    current_section_text = []
    metadata_started = False
    metadata_text = []
    
    # Process lines after the title line
    txt_parse = lines.index(song['title']) + 1
    while txt_parse < len(lines):
        line = lines[txt_parse]
        # Look ahead to detect metadata section
        # Check if we're at the end of lyrics - typically a blank line followed by authors
        if not metadata_started and txt_parse < len(lines) - 1:
            # Check if this is a blank line followed by what looks like author names
            # or if this line itself contains metadata indicators
            if ((not line.strip() and
                 txt_parse + 1 < len(lines) and
                 lines[txt_parse+1].strip() and
                 not section_pattern.match(lines[txt_parse+1]) and
                 ',' in lines[txt_parse+1]) or
                (line.strip() and metadata_pattern.match(line) and not section_pattern.match(line))):
                # Skip the blank line if that's what triggered the metadata detection
                if not line.strip():
                    txt_parse += 1
                    line = lines[txt_parse]

                # Save the last section before starting metadata
                if current_section and current_section_text:
                    song['sections'].append({
                        'name': current_section,
                        'lyrics': '\n'.join(current_section_text).strip()
                    })
                    current_section_text = []
                
                metadata_started = True
                metadata_text.append(line)
                txt_parse += 1
                continue
        
        if metadata_started:
            metadata_text.append(line)
            txt_parse += 1
            continue
            
        # Check if this line is a section header
        section_match = section_pattern.match(line)
        if section_match:
            # Save the previous section if it exists
            if current_section and current_section_text:
                song['sections'].append({
                    'name': current_section,
                    'lyrics': '\n'.join(current_section_text).strip()
                })
                current_section_text = []
            
            # Start a new section
            current_section = line.strip()
        elif line.strip() and current_section is not None:
            # Add line to current section
            current_section_text.append(line)
        
        txt_parse += 1
    
    # Add the last section if needed
    if current_section and current_section_text and not metadata_started:
        song['sections'].append({
            'name': current_section,
            'lyrics': '\n'.join(current_section_text).strip()
        })
    
    # Process metadata
    if metadata_text:
        # Extract authors - typically the first line of metadata
        if metadata_text[0].strip():
            # Check if the first line looks like authors (not starting with CCLI or ©)
            if not any(marker in metadata_text[0] for marker in ['CCLI', '©', 'Copyright']):
                song['metadata']['authors'] = metadata_text[0].strip()
        
        # Extract the copyright line (typically contains © symbol)
        copyright_line = None
        for line in metadata_text:
            if '©' in line or 'Copyright' in line:
                copyright_line = line.strip()
                break
                
        if copyright_line:
            song['metadata']['copyright'] = copyright_line
        song['metadata']['ccli_license_number'] = ccli_lic_num
    
    return song

def song_to_json(in_file):
    """
    Runs the CCLI SongSelect txt file parser and formats to JSON

    Args:
        in_file (str): CCLI SongSelect txt file to parse

    Returns:
        parsed_song_json (str): returns a JSON string that can be used with json.loads()
        song_title (str): returns the title to be used in the pptx filename
    """
    ccli_lic_num = os.getenv('CCLI_LIC_NUM')
    logging.debug(f'CCLI License Number from env var: {ccli_lic_num}')
    lyric_file_path = in_file
    parsed_song = parse_lyrics_file(lyric_file_path, ccli_lic_num)
    parsed_song_json = json.dumps(parsed_song)
    song_title = parsed_song['title']

    return parsed_song_json, song_title


def slide_formatting(presentation):
    # Set presentation to 16:9 widescreen format    
    presentation.slide_height = Inches(7.5)
    presentation.slide_width = Inches(13.33)
    # Sets a solid black background for all default layouts
    for layout in presentation.slide_layouts:
        bg_fill = layout.background.fill
        bg_fill.solid()
        bg_fill.fore_color.rgb = RGBColor(0,0,0) # black


def create_presentation_file(filename, parsed_song, tmp_save_dir):
    prs = Presentation()
    slide_formatting(prs)
    blank_layout = prs.slide_layouts[6]
    song_json = json.loads(parsed_song)

    # Calculate positioning for text box in top half of slide
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    text_box_width = Inches(12.5)
    text_box_height = Inches(3)
    left = (slide_width - text_box_width) / 2
    top = (slide_height / 4) - (text_box_height / 2)

    # Shared slide settings
    white_txt = RGBColor(255,255,255)
    font_name = 'Arial'
    align_txt_center = PP_ALIGN.CENTER

    # Create title slide with metadata
    add_title_slide = prs.slides.add_slide(blank_layout)
    title_text_box = add_title_slide.shapes.add_textbox(left, top, text_box_width, text_box_height)
    title_text_frame = title_text_box.text_frame
    title_text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    title_text_frame.clear()
    title_p = title_text_frame.paragraphs[0]
    title_p.font.name = font_name
    title_p.font.size = Pt(40)
    title_p.alignment = align_txt_center
    title_p.font.color.rgb = white_txt
    title_p.text = song_json['title']
    # Add metadata to title slide
    song_metadata = song_json['metadata']
    meta_text_box_width = Inches(12.5)
    meta_text_box_height = Inches(0.7)
    meta_text_box_margin = Inches(0.2)
    meta_text_box_top = slide_height - meta_text_box_height - meta_text_box_margin
    meta_text_box = add_title_slide.shapes.add_textbox(meta_text_box_margin, meta_text_box_top, meta_text_box_width, meta_text_box_height)
    meta_text_frame = meta_text_box.text_frame
    meta_text_frame.vertical_anchor = MSO_ANCHOR.BOTTOM
    meta_text_frame.clear()
    meta_p = meta_text_frame.paragraphs[0]
    meta_p.font.name = font_name
    meta_p.font.size = Pt(9)
    meta_p.alignment = PP_ALIGN.LEFT
    meta_p.font.color.rgb = white_txt
    meta_authors_data = song_metadata.get('authors')
    meta_authors = f'Words and music by: {meta_authors_data}'
    meta_copyright_data = song_metadata.get('copyright')
    meta_copyright = f'\n{meta_copyright_data}'
    meta_ccli_num_data = song_metadata.get('ccli_license_number')
    meta_ccli_num = f'\nUsed by permission. CCLI #: {meta_ccli_num_data}'
    meta_p.text = meta_authors + meta_copyright + meta_ccli_num

    for section in song_json['sections']:
        # Add a new blank slide
        add_slide = prs.slides.add_slide(blank_layout)

        text_box = add_slide.shapes.add_textbox(left, top, text_box_width, text_box_height)
        text_frame = text_box.text_frame
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        text_frame.clear()
        p = text_frame.paragraphs[0]
        p.font.name = font_name
        p.font.size = Pt(28)
        p.alignment = align_txt_center
        p.font.color.rgb = white_txt
        p.text = section['lyrics']

    # Save pptx file to directory
    tmp_save_filepath = os.path.join(tmp_save_dir, filename)
    logging.debug(f'tmp_save_filepath var = {tmp_save_filepath}')
    path_filename_extension = tmp_save_filepath + '.pptx'
    logging.debug(f'path_filename_extension var = {path_filename_extension}')
    prs.save(path_filename_extension)
    return path_filename_extension

def lyric_converter(txt_file_in, tmp_save_dir):
    print('[+] Importing SongSelect txt file...')
    print('[+] Parsing song text and metadata...')
    parsed_song, pptx_filename = song_to_json(txt_file_in)
    print('[+] Creating PowerPoint slides...')
    out_file = create_presentation_file(pptx_filename, parsed_song, tmp_save_dir)
    print(f'{out_file} successfully saved')
    logging.debug(parsed_song)
    return out_file
